#!/usr/bin/env python3
"""Generate Builder Digest PPTX slides from tweet JSON data.

Style: Clean white business presentation with teal (#35c1be) brand accents.
Data source: Same JSON as gen-slides.py
Output: data/output/YYYY-MM-DD/slides/slides-YYYY-MM-DD.pptx
"""

import json
import os
import re
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ===========================================
# CONFIGURATION
# ===========================================
DATA_DIR = os.path.expanduser("~/.hermes/skills/builder-digest/data/tweets")
OUTPUT_DIR = os.path.expanduser("~/.hermes/skills/builder-digest/data/output")

TODAY = "2026-04-30"
DATE_OBJ = date(2026, 4, 30)
DATE_DISPLAY = DATE_OBJ.strftime("%B %-d, %Y")

TCO_RE = re.compile(r'https?://t\.co/\S+')

# Brand colors
TEAL = RGBColor(0x35, 0xC1, 0xBE)
DARK_TEAL = RGBColor(0x0D, 0x1B, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF7, 0xF7)
MID_GRAY = RGBColor(0x9A, 0xAF, 0xAC)
DARK_TEXT = RGBColor(0x1A, 0x28, 0x27)
TEXT_SUBTLE = RGBColor(0x6B, 0x7D, 0x7A)
TEAL_LIGHT = RGBColor(0xE0, 0xF5, 0xF4)

# Layout dimensions (standard 16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN_L = Inches(1.2)
MARGIN_R = Inches(1.2)
MARGIN_T = Inches(0.6)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R


# ===========================================
# HELPERS
# ===========================================
def clean_title(text):
    if not text:
        return ""
    return TCO_RE.sub('', text).strip()

def get_title(item):
    t = item.get('title_zh', '') or ''
    if t:
        t = clean_title(t)
        if t.strip():
            return t.strip()
    t = item.get('text_zh', '') or ''
    if t:
        t = TCO_RE.sub('', t).strip()
        words = t[:30].rsplit(' ', 1)[0] if ' ' in t[:30] else t[:25]
        if words.strip():
            return words.strip() + '...'
    t = item.get('text', '') or ''
    if t:
        t = TCO_RE.sub('', t).strip()
        words = t[:30].rsplit(' ', 1)[0] if ' ' in t[:30] else t[:25]
        if words.strip():
            return words.strip() + '...'
    return ""


def add_bg(slide, color=WHITE):
    """Set slide background to a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_brand_bar(slide):
    """Add a thin teal accent bar at the top + top-left brand label."""
    # Teal accent strip at very top
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_W, Inches(0.06)
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = TEAL
    strip.line.fill.background()

    # Brand label: 牛雪梨AI
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2),
        Inches(3), Inches(0.4)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "牛雪梨AI  |  Builder Digest"
    run.font.size = Pt(10)
    run.font.color.rgb = TEXT_SUBTLE
    run.font.name = "Arial"


def add_teal_bar_left(slide, top, height):
    """Add a thin teal vertical bar on the left side."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_L - Inches(0.3), top,
        Inches(0.05), height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()


def add_footer(slide, page_num, total):
    """Add page number at bottom-right."""
    txBox = slide.shapes.add_textbox(
        SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.45),
        Inches(1.2), Inches(0.3)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{page_num:02d} / {total:02d}"
    run.font.size = Pt(9)
    run.font.color.rgb = TEXT_SUBTLE
    run.font.name = "Arial"


def add_textbox(slide, left, top, width, height, text, font_size=11,
                color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Arial", line_spacing=1.3):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    # Line spacing
    from pptx.oxml.ns import qn
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    lineSpc = pPr.makeelement(qn('a:lnSpc'), {})
    spcPct = lineSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
    lineSpc.append(spcPct)
    pPr.append(lineSpc)
    return tf


def add_multiline_textbox(slide, left, top, width, height, text, font_size=11,
                          color=DARK_TEXT, font_name="Arial", line_spacing=1.3,
                          alignment=PP_ALIGN.LEFT):
    """Add a text box with text that may contain <br> or newline separators."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Split by double newline first (paragraphs), then single newline (line breaks)
    paragraphs = text.split('\n\n')
    for pi, para_text in enumerate(paragraphs):
        if pi == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        lines = para_text.split('\n')
        for li, line_text in enumerate(lines):
            if li == 0:
                run = p.add_run()
            else:
                # Add line break
                br_elem = p._p.makeelement('{http://schemas.openxmlformats.org/drawingml/2006/main}br', {})
                p._p.append(br_elem)
                run = p.add_run()

            run.text = line_text
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = font_name

        p.space_after = Pt(4)
        p.alignment = alignment

        # Line spacing
        from pptx.oxml.ns import qn
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lineSpc = pPr.makeelement(qn('a:lnSpc'), {})
        spcPct = lineSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
        lineSpc.append(spcPct)
        pPr.append(lineSpc)

    return tf


# ===========================================
# SLIDE BUILDERS
# ===========================================
def build_title_slide(prs, total_items):
    """Build the cover slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_bg(slide, WHITE)

    # Teal accent bar at top
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_W, Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    # Teal accent bar at bottom
    bar_b = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), SLIDE_H - Inches(0.08),
        SLIDE_W, Inches(0.08)
    )
    bar_b.fill.solid()
    bar_b.fill.fore_color.rgb = TEAL
    bar_b.line.fill.background()

    # Large teal accent line left side
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2.5), Inches(2.0),
        Inches(0.08), Inches(2.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    # Title
    add_textbox(slide, Inches(3.0), Inches(2.0), Inches(8), Inches(1.0),
                "Builder Digest", font_size=44, color=DARK_TEAL, bold=True)

    # Subtitle
    add_textbox(slide, Inches(3.0), Inches(3.1), Inches(8), Inches(0.6),
                f"今日创造者精选 \u00b7 {DATE_DISPLAY}", font_size=18, color=TEXT_SUBTLE)

    # Items count
    add_textbox(slide, Inches(3.0), Inches(3.8), Inches(8), Inches(0.5),
                f"精选 {total_items} 条推文 \u00b7 源自 X/Twitter", font_size=13, color=MID_GRAY)

    # Brand label at bottom
    add_textbox(slide, Inches(3.0), Inches(5.5), Inches(8), Inches(0.5),
                "牛雪梨AI", font_size=16, color=TEAL, bold=True)


def build_content_slide(prs, item, idx, total):
    """Build a single content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    # Thin teal accent bar at top
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_W, Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    # Brand bar top-left
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(3), Inches(0.4),
                "牛雪梨AI  |  Builder Digest", font_size=9, color=TEXT_SUBTLE)

    # Page number
    add_footer(slide, idx, total)

    # Teal accent line
    add_teal_bar_left(slide, Inches(0.8), Inches(5.8))

    # --- Title ---
    title = get_title(item)
    title_y = Inches(1.0)
    add_textbox(slide, MARGIN_L + Inches(0.2), title_y, CONTENT_W - Inches(0.2), Inches(0.8),
                title, font_size=22, color=DARK_TEAL, bold=True, line_spacing=1.2)

    # --- Author info (below title) ---
    author = item.get('author_name', '') or ''
    handle = item.get('author_handle', '') or ''
    author_y = Inches(1.85)
    author_text = f"{author}  @{handle}"
    add_textbox(slide, MARGIN_L + Inches(0.2), author_y, CONTENT_W - Inches(0.2), Inches(0.35),
                author_text, font_size=11, color=TEAL, bold=False)

    # --- English text ---
    text = item.get('text', '') or ''
    text_clean = TCO_RE.sub('', text).strip()

    text_y = Inches(2.5)
    if text_clean:
        # Check if text is long
        fs = 12 if len(text_clean) > 300 else 13
        add_multiline_textbox(slide, MARGIN_L + Inches(0.2), text_y,
                              CONTENT_W - Inches(0.2), Inches(2.0),
                              text_clean, font_size=fs, color=DARK_TEXT,
                              line_spacing=1.4, font_name="Arial")

    # --- Chinese translation ---
    text_zh = item.get('text_zh', '') or ''
    text_zh_clean = TCO_RE.sub('', text_zh).strip()

    zh_y = text_y + Inches(0.2)
    if text_zh_clean:
        if text_clean:
            zh_y = text_y + Inches(2.2)
        # Check length
        fs = 11 if len(text_zh_clean) > 200 else 12
        add_multiline_textbox(slide, MARGIN_L + Inches(0.2), zh_y,
                              CONTENT_W - Inches(0.2), Inches(2.0),
                              text_zh_clean, font_size=fs, color=TEXT_SUBTLE,
                              line_spacing=1.4, font_name="Arial")

    # --- Source link at bottom ---
    url = item.get('url', '') or ''
    if url:
        add_textbox(slide, MARGIN_L + Inches(0.2), SLIDE_H - Inches(0.7),
                    CONTENT_W - Inches(0.2), Inches(0.35),
                    f"查看全文: {url}", font_size=9, color=MID_GRAY)


def build_closing_slide(prs):
    """Build the closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    # Teal accent bars top and bottom
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_W, Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    bar_b = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), SLIDE_H - Inches(0.08),
        SLIDE_W, Inches(0.08)
    )
    bar_b.fill.solid()
    bar_b.fill.fore_color.rgb = TEAL
    bar_b.line.fill.background()

    # Teal accent vertical
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2.5), Inches(2.0),
        Inches(0.08), Inches(2.0)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    # Title
    add_textbox(slide, Inches(3.0), Inches(2.0), Inches(8), Inches(0.8),
                "明天见", font_size=36, color=DARK_TEAL, bold=True)

    # Subtitle
    add_textbox(slide, Inches(3.0), Inches(2.9), Inches(8), Inches(0.5),
                "牛雪梨AI  \u00b7  Builder Digest", font_size=16, color=TEXT_SUBTLE)

    # Contact info
    contacts = [
        "订阅 & 合作",
        "",
        "X:  @niuxueli_ai",
        "Email:  915746437@qq.com",
    ]
    contact_text = "\n".join(contacts)
    add_multiline_textbox(slide, Inches(3.0), Inches(3.6), Inches(8), Inches(2.0),
                          contact_text, font_size=13, color=MID_GRAY, line_spacing=1.5)


def main():
    # Read data files
    data_path = os.path.join(DATA_DIR, TODAY)
    if not os.path.isdir(data_path):
        print(f"ERROR: Data directory not found: {data_path}")
        return 1

    files = sorted([f for f in os.listdir(data_path) if f.endswith('.json')])
    if not files:
        print(f"ERROR: No JSON files found in {data_path}")
        return 1

    items = []
    for f in files:
        with open(os.path.join(data_path, f)) as fp:
            try:
                items.append(json.load(fp))
            except json.JSONDecodeError:
                print(f"WARNING: Skipping invalid JSON: {f}")
                continue

    total_items = len(items)
    print(f"Loaded {total_items} items from {data_path}")

    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Build slides
    build_title_slide(prs, total_items)

    for idx, item in enumerate(items):
        build_content_slide(prs, item, idx + 1, total_items)

    build_closing_slide(prs)

    # Write output
    output_path = os.path.join(OUTPUT_DIR, TODAY, "slides")
    os.makedirs(output_path, exist_ok=True)
    output_file = os.path.join(output_path, f"slides-{TODAY}.pptx")

    prs.save(output_file)

    total_slides = total_items + 2
    print(f"Generated {total_slides} slides -> {output_file}")
    print(f"  - 1 title slide")
    print(f"  - {total_items} content slides")
    print(f"  - 1 closing slide")

    return 0


if __name__ == "__main__":
    exit(main())
